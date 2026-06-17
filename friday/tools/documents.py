from __future__ import annotations

from pathlib import Path

from friday.tools._decorators import register_tool


def _resolve(path: str) -> Path:
    return Path(path).expanduser().resolve()


def _normalize_slide_dict(slide: dict) -> dict[str, str]:
    """兼容 LLM 常用字段名（heading/body/content/points）。"""
    title = (
        slide.get("title")
        or slide.get("heading")
        or slide.get("name")
        or ""
    )
    bullets = (
        slide.get("bullets")
        or slide.get("body")
        or slide.get("content")
        or slide.get("text")
        or slide.get("points")
        or ""
    )
    if isinstance(bullets, list):
        bullets = "\n".join(str(item).strip() for item in bullets if str(item).strip())
    return {"title": str(title).strip(), "bullets": str(bullets).strip()}


def _slide_has_content(slide: dict[str, str]) -> bool:
    return bool(slide.get("title") or slide.get("bullets"))


def _count_pptx_text_chars(path: Path) -> int:
    from pptx import Presentation

    prs = Presentation(str(path))
    total = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            total += len(shape.text_frame.text.strip())
    return total


@register_tool(
    name="create_docx",
    description="创建 Word 文档",
    parameters={
        "type": "object",
        "properties": {
            "output_path": {"type": "string"},
            "title": {"type": "string"},
            "sections": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "heading": {"type": "string"},
                        "body": {"type": "string"},
                    },
                },
            },
        },
        "required": ["output_path", "title", "sections"],
    },
)
def create_docx(output_path: str, title: str, sections: list[dict]) -> str:
    from docx import Document

    target = _resolve(output_path)
    target.parent.mkdir(parents=True, exist_ok=True)

    doc = Document()
    doc.add_heading(title, level=0)
    for section in sections:
        heading = section.get("heading")
        body = section.get("body", "")
        if heading:
            doc.add_heading(str(heading), level=1)
        if body:
            doc.add_paragraph(str(body))
    doc.save(target)
    return f"已创建 Word 文档: {target}"


@register_tool(
    name="create_pptx",
    description="创建极简 PowerPoint 草稿（无设计、仅标题+要点）。正式汇报/演示请走内置 ppt-master 工作流，不要用本工具。",
    parameters={
        "type": "object",
        "properties": {
            "output_path": {"type": "string"},
            "title": {"type": "string"},
            "slides": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string", "description": "幻灯片标题"},
                        "bullets": {
                            "type": "string",
                            "description": "每行一条要点（也可用 heading/body，但推荐 title/bullets）",
                        },
                    },
                },
            },
        },
        "required": ["output_path", "title", "slides"],
    },
)
def create_pptx(output_path: str, title: str, slides: list[dict]) -> str:
    from pptx import Presentation
    from pptx.util import Pt

    deck_title = str(title or "").strip()
    if not deck_title:
        return "错误：title 不能为空。"

    normalized = [_normalize_slide_dict(s if isinstance(s, dict) else {}) for s in (slides or [])]
    if normalized and not any(_slide_has_content(s) for s in normalized):
        return (
            "错误：slides 中没有任何有效内容。请使用 title + bullets（每行一条要点）；"
            "勿用 heading/body（那是 create_docx 的字段名）。"
        )

    target = _resolve(output_path)
    target.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = deck_title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = "由星期五生成"

    for slide in normalized:
        if not _slide_has_content(slide):
            continue
        layout = prs.slide_layouts[1]
        page = prs.slides.add_slide(layout)
        page.shapes.title.text = slide["title"] or "要点"
        body = page.placeholders[1].text_frame
        body.clear()
        first = True
        for line in slide["bullets"].splitlines():
            if not line.strip():
                continue
            p = body.paragraphs[0] if first else body.add_paragraph()
            first = False
            p.text = line.strip()
            p.font.size = Pt(18)
            p.level = 0

    prs.save(target)
    if _count_pptx_text_chars(target) < len(deck_title):
        target.unlink(missing_ok=True)
        return (
            "错误：生成的 PPT 几乎无文字，已取消保存。"
            "请检查 slides 是否使用 title 与 bullets 字段。"
        )
    return f"已创建 PPT 文档: {target}"
