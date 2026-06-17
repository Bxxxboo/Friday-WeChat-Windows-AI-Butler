from __future__ import annotations

import tempfile
from pathlib import Path

from friday.tools.documents import _normalize_slide_dict, create_pptx


def test_normalize_slide_dict_heading_body_aliases():
    slide = _normalize_slide_dict({"heading": "RSA", "body": "公钥加密\n数字签名"})
    assert slide["title"] == "RSA"
    assert "公钥加密" in slide["bullets"]


def test_create_pptx_rejects_all_empty_slides():
    with tempfile.TemporaryDirectory() as tmp:
        out = Path(tmp) / "empty.pptx"
        result = create_pptx(str(out), "复习", [{"title": "", "bullets": ""}])
        assert result.startswith("错误：")
        assert not out.exists()


def test_create_pptx_accepts_heading_body_fields():
    with tempfile.TemporaryDirectory() as tmp:
        out = Path(tmp) / "review.pptx"
        result = create_pptx(
            str(out),
            "应用密码学期末复习",
            [{"heading": "RSA", "body": "密钥生成\n加解密流程"}],
        )
        assert result.startswith("已创建 PPT")
        assert out.exists()
        from pptx import Presentation

        prs = Presentation(str(out))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        joined = "\n".join(texts)
        assert "RSA" in joined
        assert "密钥生成" in joined


def test_create_pptx_skips_blank_content_slides():
    with tempfile.TemporaryDirectory() as tmp:
        out = Path(tmp) / "sparse.pptx"
        result = create_pptx(
            str(out),
            "标题页",
            [
                {"title": "", "bullets": ""},
                {"heading": "有内容", "body": "要点一"},
            ],
        )
        assert result.startswith("已创建 PPT")
        from pptx import Presentation

        prs = Presentation(str(out))
        assert len(prs.slides) == 2  # 封面 + 1 内容页
